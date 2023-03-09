import os
from dotenv import load_dotenv
from pptx import Presentation
from chatgptapi_translator import ChatGPTAPI
from utils import LANGUAGES, TO_LANGUAGE_CODE
import re


def get_paragraph_text(paragraph):
    """Get the text from a paragraph object"""
    text = ""
    for run in paragraph.runs:
        text += run.text
    return text


def gtp_translate(text):
    """Translate the text"""
    
    # Ignore some text
    if text.strip() == "ARiGATAYA": return text
    if text.strip() == "ARiGATAYA Entab": return text
    if len(text.strip()) == 1: return text  # Skip single characters
    
    print("GPT translating text...")
    result = translate_model.translate(text)
    return result


def replace_text(paragraph):
    """Replace the text of a paragraph object"""
    if paragraph.text.strip() == "": 
        print("Paragraph empty") 
        return
    
    if len(paragraph.runs) == 0: 
        print("Paragraph has no runs")
        return
    
    if re.findall(r'[\u4e00-\u9fff]+', paragraph.text) == [] and re.findall(r'[\u3040-\u30ff]+', paragraph.text) == []: 
        print("Paragraph has no Chinese/Japanese characters")
        return
    
    paragraph_text = get_paragraph_text(paragraph)
    print("Paragraph text: " + paragraph_text)
    
    # Process the text
    translated_text =  gtp_translate(paragraph_text)
    
    if translated_text.strip() == paragraph_text.strip(): 
        print("Skipping translation")
        print("-")
        return
    print("-")
    
    # Replace the text
    for i, run in enumerate(paragraph.runs):
        if i == 0:
            run.text = translated_text
        else:
            run.text = ""


def process_pptx_text(fileBasename):
    """Process the text in a pptx file"""
    prs = Presentation(fileBasename + '.pptx')
    
    paragraph_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:  # loop through shapes on slide
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    paragraph_count += 1
            if shape.has_table:
                for cell in shape.table.iter_cells():
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph_count += 1
    
    process_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:  # loop through shapes on slide
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    print(f"★ Processing paragraph {process_count} / {paragraph_count} for file {fileBasename}")
                    replace_text(paragraph)
                    process_count += 1
            if shape.has_table:
                for cell in shape.table.iter_cells():
                    for paragraph in cell.text_frame.paragraphs:
                        print(f"★ Processing paragraph {process_count} / {paragraph_count} for file {fileBasename}")
                        replace_text(paragraph)
                        process_count += 1
                        
    prs.save(fileBasename + '_translated.pptx')


# Create the translator
translate_model = ChatGPTAPI(
    key=os.getenv("CHATGPTAPI_KEY"),
    language=LANGUAGES.get("zh-hant"))


def translate():
    for fileBasename in os.getenv("FILE_BASENAME").split(","):
        print("=== Translate file: " + fileBasename + ".pptx ===")
        if os.path.isfile(fileBasename + ".pptx"):
            print("Start translating...")
            process_pptx_text(fileBasename)
            print("End translate.")
        else:
            print("File not found")


translate()