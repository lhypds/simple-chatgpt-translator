import os
import os.path
from dotenv import load_dotenv
from pptx import Presentation
import re


def get_paragraph_text(paragraph):
    """Get the text from a paragraph object"""
    text = ""
    for run in paragraph.runs:
        text += run.text
    return text


def process_text(paragraph, source_str, target_str):
    """Replace the text of a paragraph object"""
    if paragraph.text.strip() == "": 
        #print("Paragraph empty") 
        return
    
    if len(paragraph.runs) == 0: 
        #print("Paragraph has no runs")
        return
    
    if re.findall(r'[\u4e00-\u9fff]+', paragraph.text) == [] and re.findall(r'[\u3040-\u30ff]+', paragraph.text) == []: 
        #print("Paragraph has no Chinese/Japanese characters")
        return
    
    paragraph_text = get_paragraph_text(paragraph)
    if source_str not in paragraph_text:
        #print("Source string not found")
        return
    
    # Found text
    print("★ Paragraph text: " + paragraph_text)
    if source_str == target_str:
        #print("Source string and target string are the same")
        return
    
    result_text = paragraph_text.replace(source_str, target_str)
    print("★ Replace to text: " + result_text)
    
    # Replace the text
    for i, run in enumerate(paragraph.runs):
        if i == 0:
            run.text = result_text
        else:
            run.text = ""
    print("-")
    

def process_pptx_text(fileBasename, source_str, target_str):
    """Process the text in a pptx file"""
    prs = Presentation(fileBasename + ".pptx")
    
    for slide in prs.slides:
        for shape in slide.shapes:  # loop through shapes on slide
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    process_text(paragraph, source_str, target_str)
            if shape.has_table:
                for cell in shape.table.iter_cells():
                    for paragraph in cell.text_frame.paragraphs:
                        process_text(paragraph, source_str, target_str)
                        
    prs.save(fileBasename + ".pptx")


def replace(source_str, target_str):
    print("=== Replace " + source_str + " with " + target_str + " ===")
    for fileBasename in os.getenv("FILE_BASENAME").split(","):
        print("Replacing for file: " + fileBasename + ".pptx")
        print("Start replacing...")
        process_pptx_text(fileBasename, source_str, target_str)
        print("End replace.")
        

def search(target_str):
    print("=== Search " + target_str + " ===")
    for fileBasename in os.getenv("FILE_BASENAME").split(","):
        print("Search for file: " + fileBasename + ".pptx")
        if os.path.isfile(fileBasename + ".pptx"):
            print("Start searching...")
            process_pptx_text(fileBasename, target_str, target_str)
            print("End search.")
        else:
            print("File not found")


#replace("姿勢", "姿態")
#replace("引導", "誘導")
#replace("總站", "全站儀")
#replace("引导", "誘導")
#replace("身高", "高")
#replace("員工", "標尺")
#replace("员工", "標尺")
#replace("股份有限公司", "株式会社")
#replace("工作人員", "標尺")
#replace("普利茲姆", "棱镜")
#replace("背點", "後視點")
#replace("計測", "測量")
#replace("計量", "測量")
#replace("泰普", "類型")
#replace("人員", "標尺")
#replace("操作手冊", "操作説明書")
#replace("概要", "摘要")
#replace("切口", "掘進面")
#replace("儀器點", "機械點")
#replace("儀器", "機械")
#replace("分段", "環片")
#replace("區段", "環片")
#replace("盾構機器", "盾構機")
#replace("分節", "環片")
#replace("片段", "環片")
#replace("示威", "演示")
#replace("抗議", "演示")
#replace("導引", "誘導")
#replace("探知", "檢測")
#replace("激光", "鐳射")
#replace("縮尺", "比例尺")
#replace("導向", "誘導")
#replace("计测", "測量")
#replace("標靶", "目標")
#replace("自動校準", "自動對準")
#replace("補償量", "補償")
#replace("水准", "水平儀")
#replace("水平儀儀", "水平儀")
#replace("切割面", "掘進面")
#replace("全站仪", "全站儀")
#replace("自动", "自動")
#replace("顶", "頂")
#replace("切羽", "掘進面")
#replace("信息", "資訊")
#replace("简易", "簡易")
#replace("戒指", "環")
#replace("編號", "號")
#replace("投球", "俯仰角")
#replace("衝擊", "行程")
#replace("手續", "步驟")
#replace("標尺方式", "標尺法")
#replace("圈號", "環號")
#replace("光谱", "棱鏡")
#replace("分割片", "環片")
#replace("環號碼", "環號")
#replace("環號碼", "環號")
#replace("螢幕", "畫面")
#replace("准", "準")
#replace("鏡子", "棱鏡")