from pptx import Presentation
from chatgptapi_translator import ChatGPTAPI
from utils import LANGUAGES, TO_LANGUAGE_CODE


def get_paragraph_text(paragraph):
    """Get the text from a paragraph object"""
    text = ""
    for run in paragraph.runs:
        text += run.text
    return text


def translate_text(text):
    """Translate the text"""
    print("GPT translating text...")
    return translate_model.translate(text)


def replace_text(paragraph):
    """Replace the text of a paragraph object"""
    if paragraph.text.strip() == "": return
    if len(paragraph.runs) == 0: return
    
    paragraph_text = get_paragraph_text(paragraph)
    print("Paragraph text: " + paragraph_text)
    
    # Process the text
    for i, run in enumerate(paragraph.runs):
        if i == 0:
            run.text = translate_text(paragraph_text)
        else:
            run.text = ""


def process_pptx_text(filename):
    """Process the text in a pptx file"""
    prs = Presentation(filename)
    for slide in prs.slides:
        for shape in slide.shapes:  # loop through shapes on slide
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    replace_text(paragraph)
            if shape.has_table:
                for cell in shape.table.iter_cells():
                    for paragraph in cell.text_frame.paragraphs:
                        replace_text(paragraph)
    prs.save('test1.pptx')


# Create the translator
translate_model = ChatGPTAPI(
    key="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"
      + ",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"
      + ",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"
      + ",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"
      + ",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",
    language=LANGUAGES.get("zh-hant"))


process_pptx_text('test.pptx')
